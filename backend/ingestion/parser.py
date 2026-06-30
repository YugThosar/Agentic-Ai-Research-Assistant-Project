import os
import pandas as pd
from typing import List, Dict, Any

# Safe imports for formatting libraries
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

class DocumentParser:
    """
    Parses various document types into standard dictionary structures
    containing page/cell data and metadata.
    """
    
    @staticmethod
    def parse(file_path: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Parses a file and returns a list of pages/sections.
        Each item: {"text": str, "page_number": int, "metadata": dict}
        """
        file_type = file_type.lower()
        if file_type == "pdf":
            return DocumentParser._parse_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            return DocumentParser._parse_docx(file_path)
        elif file_type in ["pptx", "ppt"]:
            return DocumentParser._parse_pptx(file_path)
        elif file_type in ["xlsx", "xls"]:
            return DocumentParser._parse_xlsx(file_path)
        elif file_type == "csv":
            return DocumentParser._parse_csv(file_path)
        elif file_type in ["md", "markdown", "txt"]:
            return DocumentParser._parse_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    @staticmethod
    def _parse_pdf(file_path: str) -> List[Dict[str, Any]]:
        if not PdfReader:
            raise ImportError("pypdf is not installed. Run `pip install pypdf` to support PDF parsing.")
        
        pages = []
        reader = PdfReader(file_path)
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({
                "text": text.strip(),
                "page_number": idx + 1,
                "metadata": {"source": os.path.basename(file_path)}
            })
        return pages

    @staticmethod
    def _parse_docx(file_path: str) -> List[Dict[str, Any]]:
        if not docx:
            raise ImportError("python-docx is not installed. Run `pip install python-docx` to support Word document parsing.")
        
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
        
        # Word docs don't have native pagination in python-docx easily, so we treat it as one page or group paragraphs
        text_content = "\n".join(full_text)
        return [{
            "text": text_content,
            "page_number": 1,
            "metadata": {"source": os.path.basename(file_path)}
        }]

    @staticmethod
    def _parse_pptx(file_path: str) -> List[Dict[str, Any]]:
        if not Presentation:
            raise ImportError("python-pptx is not installed. Run `pip install python-pptx` to support PowerPoint parsing.")
        
        prs = Presentation(file_path)
        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slides.append({
                "text": "\n".join(slide_text),
                "page_number": idx + 1,
                "metadata": {"source": os.path.basename(file_path)}
            })
        return slides

    @staticmethod
    def _parse_xlsx(file_path: str) -> List[Dict[str, Any]]:
        sheets = pd.read_excel(file_path, sheet_name=None)
        pages = []
        for sheet_name, df in sheets.items():
            csv_like_text = df.to_csv(index=False)
            pages.append({
                "text": f"Sheet: {sheet_name}\n{csv_like_text}",
                "page_number": len(pages) + 1,
                "metadata": {"source": os.path.basename(file_path), "sheet_name": sheet_name}
            })
        return pages

    @staticmethod
    def _parse_csv(file_path: str) -> List[Dict[str, Any]]:
        df = pd.read_csv(file_path)
        csv_text = df.to_csv(index=False)
        return [{
            "text": csv_text,
            "page_number": 1,
            "metadata": {"source": os.path.basename(file_path)}
        }]

    @staticmethod
    def _parse_text(file_path: str) -> List[Dict[str, Any]]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [{
            "text": text.strip(),
            "page_number": 1,
            "metadata": {"source": os.path.basename(file_path)}
        }]
